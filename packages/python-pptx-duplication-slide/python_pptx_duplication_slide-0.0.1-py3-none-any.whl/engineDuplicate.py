from pptx import Presentation
from pptx.util import Inches
import copy, io


def MoveSlide(prs, old_index, new_index):
    xml_slides = prs.slides._sldIdLst
    slides = list(xml_slides)
    xml_slides.remove(slides[old_index])
    xml_slides.insert(new_index, slides[old_index])

def DuplicateSlide(pres, slide_index):

    # memastikan indeks slide berada dalam rentang yang benar
    if slide_index < 0 | slide_index >= len(pres.slides):
        raise IndexError("Indeks slide berada di luar rentang yang valid.")

    # ambil slide berdasarkan indeks
    slide = pres.slides[slide_index]
    slide_layout = slide.slide_layout

    # tambah slide baru dengan layout yang sama
    new_slide = pres.slides.add_slide(slide_layout)

    # copy semua shape dari slide lama ke slide baru
    for shape in slide.shapes:
        if shape.shape_type == 13:  # jika shape adalah image, jadi 13 adalah kode untuk image 
            # ambil image dari slide lama
            image_stream = io.BytesIO(shape.image.blob)
            # menambahkan image ke slide baru
            new_slide.shapes.add_picture(
                image_stream, shape.left, shape.top, shape.width, shape.height
            )
        else:
            # jika shape non-image, dilakukan deep copy seperti sebelumnya
            el = shape.element
            new_el = copy.deepcopy(el)
            new_slide.shapes._spTree.insert_element_before(new_el, 'p:extLst')

    # mengosongkan teks pada shape judul agar tampilan menjadi clear
    if new_slide.shapes.title:
        new_slide.shapes.title.text = " "

    return new_slide

# prs = Presentation('[Pemkot Yogyakarta] Report ISA Periode 02-08 Juni 2024.pptx')

# try:
#     duplicate_slide(prs, 5)

#     move_slide(prs, 29, 6)

#     prs.save('contoh_duplicated.pptx')

#     print("Slide berhasil diduplikasi dan disimpan.")
# except IndexError as e:
#     print(f"Terjadi kesalahan: {e}")
# except Exception as e:
#     print(f"Terjadi kesalahan tak terduga: {e}")
