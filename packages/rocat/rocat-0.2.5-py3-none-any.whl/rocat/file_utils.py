# rocat/file_utils.py
from openpyxl import load_workbook
import xlrd
import re
import csv
from pypdf import PdfReader
import docx2txt
import pptx
import docx2txt
from pptx import Presentation
from olefile import OleFileIO
import olefile
import zlib
import io
import zipfile
import xml.etree.ElementTree as ET
import struct



def get_txt(location, encoding='utf-8'):
    """ 
    텍스트 파일의 내용을 읽어옵니다.
    
    Args:
        location (str): 파일 경로.
        encoding (str): 파일 인코딩. 기본값은 'utf-8'입니다.
    
    Returns:
        str: 파일의 내용.
    """
    with open(location, "r", encoding=encoding) as file:
        return file.read()


def _get_column_letter(col_idx):
    col_name = ""
    while col_idx > 0:
        col_idx, remainder = divmod(col_idx - 1, 26)
        col_name = chr(65 + remainder) + col_name
    return col_name

def get_xls(file_path, sheet_index=0):
    """ 
    XLS 파일의 데이터를 딕셔너리 형태로 읽어옵니다.
    
    Args:
        file_path (str): 파일 경로.
        sheet_index (int or str): 시트 인덱스 또는 이름. 기본값은 0입니다.
    
    Returns:
        dict: 시트의 데이터가 저장된 딕셔너리.
        str: 에러 메시지.
    """
    try:
        wb = xlrd.open_workbook(file_path)
        if isinstance(sheet_index, int):
            ws = wb.sheet_by_index(sheet_index)
        else:
            ws = wb.sheet_by_name(sheet_index)

        data_dict = {}
        for col_idx in range(ws.ncols):
            col_data = ws.col_values(col_idx)
            col_name = xlrd.formula.colname(col_idx)
            data_dict[col_name] = col_data
        return data_dict
    except Exception as e:
        return f"ERROR: {str(e)}"
    

def get_xlsx(file_path, sheet=""):
    """ 
    XLSX 파일의 데이터를 딕셔너리 형태로 읽어옵니다.
    
    Args:
        file_path (str): 파일 경로.
        sheet (str): 시트 이름. 기본값은 빈 문자열로 첫 번째 시트를 읽습니다.
    
    Returns:
        dict: 시트의 데이터가 저장된 딕셔너리.
        str: 에러 메시지.
    """
    try:
        wb = load_workbook(file_path)
        if sheet == "":
            ws = wb.active
        else:
            ws = wb[sheet]
        data_dict = {}
        for col_idx, col in enumerate(ws.iter_cols(values_only=True), start=1):
            col_name = _get_column_letter(col_idx)
            data_dict[col_name] = list(col)
        return data_dict
    except Exception as e:
        return f"ERROR: {str(e)}"

def _get_column_index(col_name):
    col_idx = 0
    for char in col_name:
        col_idx = col_idx * 26 + (ord(char) - 64)
    return col_idx

def excel_get_field(field_number, excel_text):
    """ 
    엑셀 텍스트에서 지정된 필드 번호에 해당하는 데이터를 가져옵니다.
    
    Args:
        field_number (str or list): 필드 번호 또는 필드 번호의 리스트.
        excel_text (dict): 엑셀 데이터가 저장된 딕셔너리.
    
    Returns:
        list: 지정된 필드의 데이터.
        str: 에러 메시지.
    """
    try:
        if type(field_number) == str:
            if ":" in field_number:
                start_column = "".join(re.findall('[a-zA-Z]+', field_number.split(":")[0]))
                end_column = "".join(re.findall('[a-zA-Z]+', field_number.split(":")[1]))
                start_row = int("".join(re.findall('[0-9]+', field_number.split(":")[0])))
                end_row = int("".join(re.findall('[0-9]+', field_number.split(":")[1])))
                result = []
                for row in range(start_row, end_row + 1):
                    for column in range(_get_column_index(start_column), _get_column_index(end_column) + 1):
                        result.append(excel_text[_get_column_letter(column)][row - 1])
                return result
            else:
                column = "".join(re.findall('[a-zA-Z]+', field_number))
                row = int("".join(re.findall('[0-9]+', field_number)))
                return excel_text[column][row - 1]
        elif type(field_number) == list:
            result = []
            for field in field_number:
                column = "".join(re.findall('[a-zA-Z]+', field))
                row = int("".join(re.findall('[0-9]+', field)))
                result.append(excel_text[column][row - 1])
            return result
    except Exception as e:
        return f"ERROR: {str(e)}"
    

def get_doc(file_path):
    """ 
    DOC 파일의 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로.
    
    Returns:
        str: 파일의 텍스트 내용.
    """
    ole = olefile.OleFileIO(file_path)
    word_doc_stream = None
    for stream in ole.listdir():
        if stream[-1].startswith('WordDocument'):
            word_doc_stream = stream
            break
    
    if not word_doc_stream:
        return "ERROR: WordDocument stream not found."
        
    data = ole.openstream(word_doc_stream).read()
    
    text = extract_doc_text(data)
    
    return text

def extract_doc_text(data):        
    fib = get_fib(data)
    
    text_start = fib['fcClx']
    text_size = fib['lcbClx']     
    rawtext = data[text_start:text_start+text_size]
    text = rawtext.decode('utf-16').strip()
    
    return text
    
def get_fib(data):
    fib = {}
    fib['wIdent'] = int.from_bytes(data[0:2], byteorder='little') 
    fib['nFib'] = int.from_bytes(data[2:4], byteorder='little')
    fib['unused'] = int.from_bytes(data[4:8], byteorder='little')
    fib['lid'] = int.from_bytes(data[8:12], byteorder='little')
    fib['pnNext'] = int.from_bytes(data[12:14], byteorder='little')
    fib['flags'] = bin(int.from_bytes(data[14:16], byteorder='little'))[2:]
    fib['nFibBack'] = int.from_bytes(data[16:18], byteorder='little')
    
    if fib['nFib'] == 193:
        fib['fcClx'] = int.from_bytes(data[68:72], byteorder='little')  
        fib['lcbClx'] = int.from_bytes(data[72:76], byteorder='little')
    elif fib['nFib'] == 217: 
        fib['fcClx'] = int.from_bytes(data[88:92], byteorder='little')
        fib['lcbClx'] = int.from_bytes(data[92:96], byteorder='little')
    else:
        return None
        
    return fib
    

def get_docx(file_path):
    try:
        text = docx2txt.process(file_path)
        return text
    except Exception as e:
        return f"ERROR: {str(e)}"
    


def get_ppt(file_path):
    pass


def get_pptx(file_path):
    """ 
    PPTX 파일의 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로.
    
    Returns:
        str: 파일의 텍스트 내용.
        str: 에러 메시지.
    """
    try:
        text = ""
        ppt = pptx.Presentation(file_path)
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"ERROR: {str(e)}"

def get_csv(file_path, delimiter=',', encoding='utf-8'):
    """ 
    CSV 파일의 데이터를 읽어옵니다.
    
    Args:
        file_path (str): 파일 경로.
        delimiter (str): CSV 파일의 구분자. 기본값은 ','입니다.
        encoding (str): 파일 인코딩. 기본값은 'utf-8'입니다.
    
    Returns:
        list: CSV 파일의 행 데이터를 저장한 리스트.
        str: 에러 메시지.
    """
    try:
        with open(file_path, 'r', encoding=encoding) as file:
            csv_reader = csv.reader(file, delimiter=delimiter)
            rows = [row for row in csv_reader]
            return rows
    except Exception as e:
        return f"ERROR: {str(e)}"
    
def get_pdf(file_path):
    """ 
    PDF 파일의 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로.
    
    Returns:
        str: PDF 파일의 텍스트 내용.
        str: 에러 메시지.
    """
    try:
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text()
            return text
    except Exception as e:
        return f"ERROR: {str(e)}"
    

def get_hwp(file_path):
    """ 
    HWP 파일의 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로.
    
    Returns:
        str: HWP 파일의 텍스트 내용.
        str: 에러 메시지.
    """
    try:
        return read_hwp(file_path)
    except Exception as e:
        return f"ERROR: {str(e)}"

def get_hwpx(file_path):
    """ 
    HWPX 파일의 텍스트를 추출합니다.
    
    Args:
        file_path (str): 파일 경로.
    
    Returns:
        str: HWPX 파일의 텍스트 내용.
        str: 에러 메시지.
    """
    try:
        return read_hwpx(file_path)
    except Exception as e:
        return f"ERROR: {str(e)}"

def read_hwp(hwp_path):
    with OleFileIO(hwp_path) as ole:
        validate_hwp_file(ole)
        compression_flag = get_compression_flag(ole)
        section_texts = [read_section(ole, section_id, compression_flag) for section_id in get_section_ids(ole)]
    
    return '\n'.join(section_texts).strip()


def validate_hwp_file(ole):
    required_streams = {"FileHeader", "\x05HwpSummaryInformation"}
    if not required_streams.issubset(set('/'.join(stream) for stream in ole.listdir())):
        raise ValueError("The file is not a valid HWP document.")

def get_compression_flag(ole):
    with ole.openstream("FileHeader") as header_stream:
        return bool(header_stream.read(37)[36] & 1)

def get_section_ids(ole):
    return sorted(
        int(stream[1].replace("Section", "")) 
        for stream in ole.listdir() if stream[0] == "BodyText"
    )

def read_section(ole, section_id, is_compressed):
    with ole.openstream(f"BodyText/Section{section_id}") as section_stream:
        data = section_stream.read()
        if is_compressed:
            data = zlib.decompress(data, -15)
        return extract_text(data)

def extract_text(data):
    text, cursor = "", 0
    while cursor < len(data):
        header = struct.unpack_from("<I", data, cursor)[0]
        type, length = header & 0x3ff, (header >> 20) & 0xfff
        if type == 67:
            text += data[cursor + 4:cursor + 4 + length].decode('utf-16') + "\n"
        cursor += 4 + length
    return text


def read_hwpx(hwpx_file_path):
    with open(hwpx_file_path, 'rb') as f:
        hwpx_file_bytes = f.read()

    with io.BytesIO(hwpx_file_bytes) as bytes_io:
        with zipfile.ZipFile(bytes_io, 'r') as zip_ref:
            text_parts = []
            for file_info in zip_ref.infolist():
                if file_info.filename.startswith('Contents/') and file_info.filename.endswith('.xml'):
                    with zip_ref.open(file_info) as file:
                        tree = ET.parse(file)
                        root = tree.getroot()
                        for elem in root.iter():
                            if elem.text:
                                text_parts.append(elem.text.strip())
    
   
    return '\n'.join(text_parts)